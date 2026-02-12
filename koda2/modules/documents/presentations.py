"""Presentation generation module for creating PowerPoint slides."""

from __future__ import annotations

from dataclasses import dataclass, field
from enum import StrEnum
from pathlib import Path
from typing import Any, Optional

from koda2.logging_config import get_logger

logger = get_logger(__name__)


class SlideLayout(StrEnum):
    """Available slide layouts."""
    TITLE = "title"
    TITLE_AND_CONTENT = "title_and_content"
    TWO_CONTENT = "two_content"
    BLANK = "blank"
    TITLE_ONLY = "title_only"
    SECTION_HEADER = "section_header"
    COMPARISON = "comparison"
    CONTENT_WITH_CAPTION = "content_with_caption"
    PICTURE_WITH_CAPTION = "picture_with_caption"


@dataclass
class SlideContent:
    """Content for a single slide."""
    layout: SlideLayout = SlideLayout.TITLE_AND_CONTENT
    title: str = ""
    content: str = ""  # Can be text, bullet points, or HTML-like markup
    subtitle: str = ""
    notes: str = ""  # Speaker notes
    images: list[str] = field(default_factory=list)  # Paths to images
    chart_data: Optional[dict] = None  # For charts
    table_data: Optional[list[list[str]]] = None  # For tables
    background_image: Optional[str] = None


@dataclass
class PresentationTheme:
    """Theme configuration for presentations."""
    name: str = "Corporate"
    primary_color: str = "1F4E78"  # RGB hex
    secondary_color: str = "2E75B6"
    accent_color: str = "5B9BD5"
    text_color: str = "333333"
    background_color: str = "FFFFFF"
    font_title: str = "Calibri"
    font_body: str = "Calibri"
    logo_path: Optional[str] = None
    footer_text: str = ""


class PresentationGenerator:
    """Generate professional PowerPoint presentations."""
    
    def __init__(self, theme: Optional[PresentationTheme] = None):
        self.theme = theme or PresentationTheme()
        
    def create_presentation(
        self,
        title: str,
        subtitle: str = "",
        author: str = "",
        slides: Optional[list[SlideContent]] = None,
        output_path: Optional[str] = None,
    ) -> str:
        """Create a complete presentation."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        
        prs = Presentation()
        prs.core_properties.title = title
        prs.core_properties.author = author
        prs.core_properties.subject = subtitle
        
        # Set slide dimensions (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Add title slide
        self._add_title_slide(prs, title, subtitle, author)
        
        # Add content slides
        if slides:
            for slide_content in slides:
                self._add_slide(prs, slide_content)
                
        # Add closing slide
        self._add_closing_slide(prs)
        
        # Save
        if not output_path:
            output_path = f"data/presentations/{title.replace(' ', '_')}.pptx"
            
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        prs.save(str(output_path))
        logger.info("presentation_created", path=str(output_path), slides=len(slides or []) + 2)
        return str(output_path)
        
    def _add_title_slide(
        self, 
        prs: Any, 
        title: str, 
        subtitle: str,
        author: str,
    ) -> None:
        """Add title slide."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Set subtitle
        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_frame = subtitle_shape.text_frame
            subtitle_frame.paragraphs[0].font.size = Pt(24)
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
        # Add author/date at bottom
        if author:
            left = Inches(1)
            top = Inches(6.5)
            width = Inches(11)
            height = Inches(0.5)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = f"{author} | {__import__('datetime').date.today().strftime('%B %Y')}"
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.paragraphs[0].font.size = Pt(14)
            
    def _add_slide(self, prs: Any, content: SlideContent) -> None:
        """Add a content slide based on layout."""
        layout_methods = {
            SlideLayout.TITLE: self._layout_title,
            SlideLayout.TITLE_AND_CONTENT: self._layout_title_and_content,
            SlideLayout.TWO_CONTENT: self._layout_two_content,
            SlideLayout.BLANK: self._layout_blank,
            SlideLayout.TITLE_ONLY: self._layout_title_only,
            SlideLayout.SECTION_HEADER: self._layout_section_header,
        }
        
        method = layout_methods.get(content.layout, self._layout_title_and_content)
        method(prs, content)
        
    def _layout_title_and_content(self, prs: Any, content: SlideContent) -> None:
        """Title and content layout."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        slide.shapes.title.text = content.title
        
        # Content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = content.content
        
        # Format bullet points if content has multiple lines
        lines = content.content.split('\n')
        if len(lines) > 1:
            tf.clear()
            for i, line in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line.lstrip('- ').strip()
                p.level = 0 if not line.startswith('  ') else 1
                p.font.size = Pt(18)
                
        # Add table if present
        if content.table_data:
            self._add_table_to_slide(slide, content.table_data)
            
        # Add notes
        if content.notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = content.notes
            
    def _layout_title_only(self, prs: Any, content: SlideContent) -> None:
        """Title only layout."""
        slide_layout = prs.slide_layouts[5]  # Title Only
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = content.title
        
    def _layout_section_header(self, prs: Any, content: SlideContent) -> None:
        """Section header layout."""
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN
        
        slide_layout = prs.slide_layouts[2]  # Section Header
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = content.title
        if content.subtitle:
            slide.placeholders[1].text = content.subtitle
            
    def _layout_two_content(self, prs: Any, content: SlideContent) -> None:
        """Two content layout."""
        slide_layout = prs.slide_layouts[5]  # Blank, we'll add manually
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title = slide.shapes.add_textbox(
            left=prs.slide_width * 0.05,
            top=prs.slide_height * 0.05,
            width=prs.slide_width * 0.9,
            height=prs.slide_height * 0.15,
        )
        title.text_frame.text = content.title
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        
    def _layout_title(self, prs: Any, content: SlideContent) -> None:
        """Centered title layout."""
        self._layout_title_only(prs, content)
        
    def _layout_blank(self, prs: Any, content: SlideContent) -> None:
        """Blank layout."""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
    def _add_table_to_slide(self, slide: Any, data: list[list[str]]) -> None:
        """Add a table to a slide."""
        from pptx.util import Inches
        from pptx.enum.text import PP_ALIGN
        
        if not data:
            return
            
        rows = len(data)
        cols = len(data[0]) if data else 0
        
        left = Inches(1)
        top = Inches(3)
        width = Inches(11)
        height = Inches(4)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set column widths
        for i in range(cols):
            table.columns[i].width = Inches(11 / cols)
            
        # Fill data
        for i, row_data in enumerate(data):
            for j, cell_data in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(cell_data)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Header row formatting
                if i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(31, 78, 120)  # Dark blue
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    cell.text_frame.paragraphs[0].font.bold = True
                    
    def _add_closing_slide(self, prs: Any) -> None:
        """Add a closing 'Thank You' slide."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Add centered thank you text
        left = Inches(2)
        top = Inches(3)
        width = Inches(9)
        height = Inches(2)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.text = "Thank You"
        tf.paragraphs[0].font.size = Pt(54)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = "Questions?"
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER
        
    def generate_from_outline(
        self,
        outline: str,
        title: str,
        output_path: Optional[str] = None,
    ) -> str:
        """Generate presentation from text outline.
        
        Outline format:
        # Main Title
        ## Slide 1 Title
        - Bullet point 1
        - Bullet point 2
        ## Slide 2 Title
        Content text
        """
        slides = []
        current_slide = None
        content_buffer = []
        
        for line in outline.strip().split('\n'):
            line = line.strip()
            if not line:
                continue
                
            if line.startswith('# '):
                # Main title - ignore, already have it
                pass
            elif line.startswith('## '):
                # Save previous slide
                if current_slide:
                    current_slide.content = '\n'.join(content_buffer)
                    slides.append(current_slide)
                    content_buffer = []
                    
                # New slide
                current_slide = SlideContent(
                    title=line[3:],
                    layout=SlideLayout.TITLE_AND_CONTENT,
                )
            elif line.startswith('- '):
                content_buffer.append(line)
            elif current_slide:
                content_buffer.append(line)
                
        # Save last slide
        if current_slide:
            current_slide.content = '\n'.join(content_buffer)
            slides.append(current_slide)
            
        return self.create_presentation(
            title=title,
            slides=slides,
            output_path=output_path,
        )


from pptx.dml.color import RGBColor
