"""Document analyzer service for extracting content from various file types."""

from __future__ import annotations

import csv
import io
from pathlib import Path
from typing import Optional, Any

from koda2.logging_config import get_logger
from koda2.modules.document_analyzer.models import DocumentAnalysis, FileType

logger = get_logger(__name__)


class DocumentAnalyzerService:
    """Service for analyzing documents and extracting structured information.
    
    Supports:
    - PDF files (text extraction)
    - Office documents (DOCX, XLSX, PPTX)
    - Text files (TXT, CSV)
    - Images (with OCR and vision analysis)
    """
    
    def __init__(self, llm_router: Optional[Any] = None) -> None:
        self._llm = llm_router
    
    async def analyze_file(self, file_path: str) -> DocumentAnalysis:
        """Analyze a file and extract its content.
        
        Args:
            file_path: Path to the file to analyze
            
        Returns:
            DocumentAnalysis with extracted content
        """
        path = Path(file_path)
        
        if not path.exists():
            return DocumentAnalysis(
                file_path=file_path,
                file_type=FileType.UNKNOWN,
                file_size=0,
                filename=path.name,
                analysis_error="File not found",
            )
        
        file_type = self._detect_file_type(path)
        file_size = path.stat().st_size
        
        try:
            if file_type == FileType.PDF:
                return await self._analyze_pdf(path)
            elif file_type == FileType.DOCX:
                return await self._analyze_docx(path)
            elif file_type == FileType.XLSX:
                return await self._analyze_xlsx(path)
            elif file_type == FileType.PPTX:
                return await self._analyze_pptx(path)
            elif file_type == FileType.IMAGE:
                return await self._analyze_image(path)
            elif file_type in (FileType.TXT, FileType.CSV):
                return await self._analyze_text(path, file_type)
            else:
                return DocumentAnalysis(
                    file_path=file_path,
                    file_type=file_type,
                    file_size=file_size,
                    filename=path.name,
                    analysis_error=f"Unsupported file type: {file_type.value}",
                )
        except Exception as exc:
            logger.error("document_analysis_failed", path=file_path, error=str(exc))
            return DocumentAnalysis(
                file_path=file_path,
                file_type=file_type,
                file_size=file_size,
                filename=path.name,
                analysis_error=str(exc),
            )
    
    def _detect_file_type(self, path: Path) -> FileType:
        """Detect file type from extension."""
        ext = path.suffix.lower()
        
        mapping = {
            ".pdf": FileType.PDF,
            ".docx": FileType.DOCX,
            ".xlsx": FileType.XLSX,
            ".pptx": FileType.PPTX,
            ".txt": FileType.TXT,
            ".csv": FileType.CSV,
            ".jpg": FileType.IMAGE,
            ".jpeg": FileType.IMAGE,
            ".png": FileType.IMAGE,
            ".gif": FileType.IMAGE,
            ".webp": FileType.IMAGE,
            ".bmp": FileType.IMAGE,
        }
        
        return mapping.get(ext, FileType.UNKNOWN)
    
    async def _analyze_pdf(self, path: Path) -> DocumentAnalysis:
        """Extract content from PDF file."""
        try:
            import pdfplumber
        except ImportError:
            logger.warning("pdfplumber_not_installed", path=str(path))
            # Fallback to basic info
            return DocumentAnalysis(
                file_path=str(path),
                file_type=FileType.PDF,
                file_size=path.stat().st_size,
                filename=path.name,
                analysis_error="PDF analysis requires pdfplumber: pip install pdfplumber",
            )
        
        text_parts = []
        metadata = {}
        
        with pdfplumber.open(path) as pdf:
            # Extract metadata
            if pdf.metadata:
                metadata = {
                    "title": pdf.metadata.get("Title"),
                    "author": pdf.metadata.get("Author"),
                    "created": pdf.metadata.get("CreationDate"),
                }
            
            # Extract text from all pages
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
        
        full_text = "\n\n".join(text_parts)
        
        # Generate AI summary if LLM available
        summary = None
        topics = []
        actions = []
        if self._llm and full_text:
            summary, topics, actions = await self._generate_summary(full_text, "PDF document")
        
        return DocumentAnalysis(
            file_path=str(path),
            file_type=FileType.PDF,
            file_size=path.stat().st_size,
            filename=path.name,
            text_content=full_text[:10000] if full_text else None,  # Limit storage
            summary=summary,
            title=metadata.get("title"),
            author=metadata.get("author"),
            key_topics=topics,
            action_items=actions,
        )
    
    async def _analyze_docx(self, path: Path) -> DocumentAnalysis:
        """Extract content from Word document."""
        try:
            from docx import Document
        except ImportError:
            return DocumentAnalysis(
                file_path=str(path),
                file_type=FileType.DOCX,
                file_size=path.stat().st_size,
                filename=path.name,
                analysis_error="DOCX analysis requires python-docx",
            )
        
        doc = Document(path)
        
        # Extract text from paragraphs
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    text_parts.append(" | ".join(row_text))
        
        full_text = "\n".join(text_parts)
        
        # Get core properties
        title = None
        author = None
        if doc.core_properties:
            title = doc.core_properties.title
            author = doc.core_properties.author
        
        # Generate AI summary
        summary = None
        topics = []
        actions = []
        if self._llm and full_text:
            summary, topics, actions = await self._generate_summary(full_text, "Word document")
        
        return DocumentAnalysis(
            file_path=str(path),
            file_type=FileType.DOCX,
            file_size=path.stat().st_size,
            filename=path.name,
            text_content=full_text[:10000] if full_text else None,
            summary=summary,
            title=title,
            author=author,
            key_topics=topics,
            action_items=actions,
        )
    
    async def _analyze_xlsx(self, path: Path) -> DocumentAnalysis:
        """Extract content from Excel spreadsheet."""
        try:
            from openpyxl import load_workbook
        except ImportError:
            return DocumentAnalysis(
                file_path=str(path),
                file_type=FileType.XLSX,
                file_size=path.stat().st_size,
                filename=path.name,
                analysis_error="XLSX analysis requires openpyxl",
            )
        
        wb = load_workbook(path, data_only=True)
        
        sheet_data = []
        total_rows = 0
        total_cols = 0
        
        for sheet_name in wb.sheetnames[:3]:  # Limit to first 3 sheets
            ws = wb[sheet_name]
            sheet_rows = []
            
            # Get header row
            headers = []
            for cell in ws[1]:
                headers.append(str(cell.value) if cell.value else "")
            
            if headers:
                sheet_rows.append(f"Sheet '{sheet_name}' columns: {', '.join(headers)}")
            
            # Sample data rows (first 10)
            for i, row in enumerate(ws.iter_rows(min_row=2, max_row=11, values_only=True), 2):
                row_values = [str(v) if v is not None else "" for v in row]
                if any(row_values):
                    sheet_rows.append(f"Row {i}: {' | '.join(row_values[:5])}")  # First 5 cols
            
            if ws.max_row > 11:
                sheet_rows.append(f"... and {ws.max_row - 11} more rows")
            
            sheet_data.append("\n".join(sheet_rows))
            total_rows += ws.max_row
            total_cols = max(total_cols, ws.max_column)
        
        full_text = f"Excel file with {len(wb.sheetnames)} sheets:\n\n"
        full_text += "\n\n---\n\n".join(sheet_data)
        
        # Generate summary
        summary = None
        topics = []
        if self._llm:
            summary, topics, _ = await self._generate_summary(full_text, "Excel spreadsheet")
        
        return DocumentAnalysis(
            file_path=str(path),
            file_type=FileType.XLSX,
            file_size=path.stat().st_size,
            filename=path.name,
            text_content=full_text[:8000] if full_text else None,
            summary=summary,
            sheet_names=wb.sheetnames,
            row_count=total_rows,
            column_count=total_cols,
            key_topics=topics,
        )
    
    async def _analyze_pptx(self, path: Path) -> DocumentAnalysis:
        """Extract content from PowerPoint presentation."""
        try:
            from pptx import Presentation
        except ImportError:
            return DocumentAnalysis(
                file_path=str(path),
                file_type=FileType.PPTX,
                file_size=path.stat().st_size,
                filename=path.name,
                analysis_error="PPTX analysis requires python-pptx",
            )
        
        prs = Presentation(path)
        
        slide_texts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_content = [f"Slide {i}:"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
            
            if len(slide_content) > 1:
                slide_texts.append("\n".join(slide_content))
        
        full_text = "\n\n".join(slide_texts)
        
        # Generate summary
        summary = None
        topics = []
        if self._llm:
            summary, topics, _ = await self._generate_summary(
                full_text, f"PowerPoint presentation ({len(prs.slides)} slides)"
            )
        
        return DocumentAnalysis(
            file_path=str(path),
            file_type=FileType.PPTX,
            file_size=path.stat().st_size,
            filename=path.name,
            text_content=full_text[:8000] if full_text else None,
            summary=summary,
            slide_count=len(prs.slides),
            key_topics=topics,
        )
    
    async def _analyze_image(self, path: Path) -> DocumentAnalysis:
        """Analyze image using vision capabilities."""
        # Use the existing ImageService for analysis
        try:
            from koda2.modules.images import ImageService
            
            image_service = ImageService()
            
            # Generate description
            description = await image_service.analyze(
                str(path),
                prompt="Describe this image in detail. If it contains text, transcribe all visible text.",
            )
            
            # Try OCR for text extraction
            detected_text = None
            try:
                import pytesseract
                from PIL import Image
                
                img = Image.open(path)
                detected_text = pytesseract.image_to_string(img)
                if not detected_text.strip():
                    detected_text = None
            except Exception:
                # OCR not available or failed, use vision description
                pass
            
            # Extract topics with LLM
            topics = []
            if self._llm and description:
                _, topics, _ = await self._generate_summary(description, "Image")
            
            return DocumentAnalysis(
                file_path=str(path),
                file_type=FileType.IMAGE,
                file_size=path.stat().st_size,
                filename=path.name,
                image_description=description,
                detected_text=detected_text,
                key_topics=topics,
            )
            
        except Exception as exc:
            logger.error("image_analysis_failed", path=str(path), error=str(exc))
            return DocumentAnalysis(
                file_path=str(path),
                file_type=FileType.IMAGE,
                file_size=path.stat().st_size,
                filename=path.name,
                analysis_error=str(exc),
            )
    
    async def _analyze_text(self, path: Path, file_type: FileType) -> DocumentAnalysis:
        """Analyze plain text or CSV file."""
        content = path.read_text(encoding="utf-8", errors="replace")
        
        if file_type == FileType.CSV:
            # Parse CSV to get structure
            lines = content.strip().split("\n")
            reader = csv.reader(lines[:21])  # First 20 rows + header
            rows = list(reader)
            
            if rows:
                preview = []
                preview.append(f"Columns: {', '.join(rows[0])}")
                for row in rows[1:6]:
                    preview.append(f"Row: {' | '.join(row[:5])}")
                if len(rows) > 6:
                    preview.append(f"... and {len(rows) - 6} more rows")
                
                content = "\n".join(preview)
        
        # Generate summary
        summary = None
        topics = []
        actions = []
        if self._llm and content:
            summary, topics, actions = await self._generate_summary(content, "Text document")
        
        return DocumentAnalysis(
            file_path=str(path),
            file_type=file_type,
            file_size=path.stat().st_size,
            filename=path.name,
            text_content=content[:10000] if content else None,
            summary=summary,
            key_topics=topics,
            action_items=actions,
        )
    
    async def _generate_summary(
        self,
        text: str,
        doc_type: str,
    ) -> tuple[Optional[str], list[str], list[str]]:
        """Generate AI summary of document content.
        
        Returns:
            Tuple of (summary, key_topics, action_items)
        """
        if not self._llm or not text or len(text.strip()) < 50:
            return None, [], []
        
        # Truncate text if too long
        max_chars = 4000
        truncated = text[:max_chars]
        if len(text) > max_chars:
            truncated += f"\n\n[Document truncated, {len(text) - max_chars} more characters...]"
        
        prompt = f"""Analyze this {doc_type} and provide:
1. A brief summary (2-3 sentences)
2. Key topics/themes (max 5)
3. Any action items or tasks mentioned (max 5)

Document content:
{truncated}

Respond in this format:
SUMMARY: <summary>
TOPICS: <comma-separated topics>
ACTIONS: <comma-separated action items, or "None">
"""
        
        try:
            from koda2.modules.llm.models import ChatMessage, LLMRequest
            
            request = LLMRequest(
                messages=[ChatMessage(role="user", content=prompt)],
                temperature=0.3,
            )
            
            response = await self._llm.complete(request)
            content = response.content
            
            # Parse response
            summary = None
            topics = []
            actions = []
            
            for line in content.split("\n"):
                line = line.strip()
                if line.startswith("SUMMARY:"):
                    summary = line[8:].strip()
                elif line.startswith("TOPICS:"):
                    topics_text = line[7:].strip()
                    topics = [t.strip() for t in topics_text.split(",") if t.strip()]
                elif line.startswith("ACTIONS:"):
                    actions_text = line[8:].strip()
                    if actions_text.lower() != "none":
                        actions = [a.strip() for a in actions_text.split(",") if a.strip()]
            
            return summary, topics, actions
            
        except Exception as exc:
            logger.error("summary_generation_failed", error=str(exc))
            return None, [], []
    
    async def analyze_with_context(
        self,
        file_path: str,
        user_message: str,
    ) -> DocumentAnalysis:
        """Analyze a file with the context of the user's message.
        
        This helps understand what the user wants to know about the file.
        """
        # First do standard analysis
        analysis = await self.analyze_file(file_path)
        
        # If we have an LLM, ask a specific question based on user message
        if self._llm and analysis.is_successful():
            try:
                context = analysis.to_context_string(max_length=3000)
                
                prompt = f"""The user sent a file with this message: "{user_message}"

Here's what I found in the file:
{context}

Based on the user's message and the file content, what would be most helpful to tell them?
Provide a concise, helpful response that addresses their likely intent.
"""
                
                from koda2.modules.llm.models import ChatMessage, LLMRequest
                
                request = LLMRequest(
                    messages=[ChatMessage(role="user", content=prompt)],
                    temperature=0.5,
                )
                
                response = await self._llm.complete(request)
                
                # Enhance the analysis with context-aware insights
                analysis.summary = response.content
                
            except Exception as exc:
                logger.error("context_analysis_failed", error=str(exc))
        
        return analysis
